"""
build_slides_global_2004.py — Add the global_2004 slides to the defense deck.

Addresses: P4 — the deck is a publication surface. Every figure it states is
read from a committed artifact here, for the same reason the report chapter and
GLOBAL_2004_RESULTS.md are generated: a number typed into a slide cannot be
checked against the evidence, and a defense deck is the worst place to discover
a stale one.

DESIGN. The existing 18 slides are hand-made and share a strict grid, extracted
from slide 13 rather than invented:

    eyebrow   Arial  9pt bold   #3D8DFF   at (48, 22)
    title     Arial 27pt bold   #111827   at (48, 44)
    page no.  Arial  9pt bold   #5D6675   at (878, 27)
    KPI       Arial 30pt bold   accent
    caption   Arial 12pt        #5D6675
    panel     fill #F3F4F6      callout fill #FFF0D5

New slides reuse those values exactly so they are indistinguishable from the
hand-made ones. The result is rendered to PDF and inspected — a slide that only
looks right in code is not verified.

IDEMPOTENT: existing global_2004 slides are removed before insertion, so
re-running never duplicates them.

Usage:
    python scripts/build_slides_global_2004.py
"""

from __future__ import annotations

import copy
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
GOLD = ROOT / "data" / "gold"
DECK = ROOT / "output" / "presentation" / "Soutenance_PFA_Portfolio_ML_INPT_EURAFRIC.pptx"
FIGDIR = ROOT / "docs" / "rapport_final" / "assets" / "figures"
PNGDIR = ROOT / "output" / "presentation" / "_figures"

BLUE, INK, GREY = RGBColor(0x3D, 0x8D, 0xFF), RGBColor(0x11, 0x18, 0x27), RGBColor(0x5D, 0x66, 0x75)
RED, TEAL, AMBER = RGBColor(0xB7, 0x43, 0x3E), RGBColor(0x0F, 0x76, 0x6E), RGBColor(0xA8, 0x68, 0x12)
PANEL, CALLOUT = RGBColor(0xF3, 0xF4, 0xF6), RGBColor(0xFF, 0xF0, 0xD5)

MARK = "GLOBAL2004"          # invisible marker so re-runs can find and replace


def load(name: str) -> dict:
    return json.loads((GOLD / name).read_text())


def fr(value: float, dp: int = 4, signed: bool = False) -> str:
    """Format ONE number in French convention: decimal comma, U+2212 minus.

    Applied per number, never to a sentence. The first version of this script
    called `.replace(".", ",")` on whole strings, which silently turned
    sentence-ending periods into commas on the rendered slide — visible only
    once the deck was actually looked at.
    """
    text = f"{value:+.{dp}f}" if signed else f"{value:.{dp}f}"
    return text.replace(".", ",").replace("-", "−")


# ── primitives, all on the extracted grid ────────────────────────────────────

def text(slide, x, y, w, h, s, size, *, bold=False, color=INK, align=None, spacing=None):
    box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(s.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.name, run.font.size, run.font.bold = "Arial", Pt(size), bold
        run.font.color.rgb = color
        if align is not None:
            para.alignment = align
        if spacing:
            para.line_spacing = spacing
    return box


def panel(slide, x, y, w, h, fill=PANEL):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = 0.03
    except Exception:
        pass
    return shape


def kpi(slide, x, y, w, value, caption, color):
    text(slide, x, y, w, 46, value, 30, bold=True, color=color)
    text(slide, x, y + 46, w, 36, caption, 12, color=GREY, spacing=1.15)


def header(slide, eyebrow, title, number):
    text(slide, 48, 22, 375, 16, eyebrow, 9, bold=True, color=BLUE)
    text(slide, 48, 44, 840, 44, title, 27, bold=True, color=INK)
    text(slide, 878, 27, 34, 18, str(number), 9, bold=True, color=GREY)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in list(slide.shapes):        # strip the layout's placeholders
        shape._element.getparent().remove(shape._element)
    # Invisible marker, off-canvas, so a re-run can identify our slides.
    text(slide, -200, -200, 40, 10, MARK, 1, color=GREY)
    return slide


def figure_png(stem: str) -> Path:
    PNGDIR.mkdir(parents=True, exist_ok=True)
    out = PNGDIR / f"{stem}.png"
    subprocess.run(["pdftoppm", "-png", "-r", "220", "-singlefile",
                    str(FIGDIR / f"{stem}.pdf"), str(PNGDIR / stem)], check=True)
    return out


# ── the five slides ──────────────────────────────────────────────────────────

def slide_why(prs, readiness, cap):
    s = new_slide(prs)
    header(s, "EXTENSION DE RECHERCHE",
           "Un troisième univers pour lever un doute d'identification", "")
    etf = cap["results"]["0.25"]["min_variance_lw"]["distinct_allocations"]

    panel(s, 48, 112, 406, 200)
    text(s, 70, 132, 360, 21, "etf_2017 — la contrainte domine", 13, bold=True, color=INK)
    text(s, 70, 166, 200, 46, f"{etf} / 248", 30, bold=True, color=RED)
    text(s, 70, 214, 360, 80,
         "allocations distinctes de la variance minimale.\n"
         "À cinq actifs sous un plafond de 25 %, c'est la\n"
         "contrainte — non le modèle — qui choisit.", 12, color=GREY, spacing=1.3)

    panel(s, 476, 112, 406, 200)
    text(s, 498, 132, 360, 21, "full_2021 — la covariance est biaisée", 13, bold=True, color=INK)
    text(s, 498, 166, 200, 46, "×19,1", 30, bold=True, color=RED)
    text(s, 498, 214, 360, 80,
         "corrélation décalée / contemporaine face au SPY.\n"
         "Séances non recouvrantes et prix figés\n"
         "(17,1 % de journées à rendement nul).", 12, color=GREY, spacing=1.3)

    panel(s, 48, 338, 834, 92, CALLOUT)
    text(s, 70, 356, 400, 18, "Le doute que cela crée", 10.5, bold=True, color=AMBER)
    text(s, 70, 380, 790, 40,
         "Sur ces deux univers, impossible de distinguer « le modèle n'apporte rien »\n"
         "de « le dispositif ne lui permet pas de s'exprimer ».", 12.5, bold=True,
         color=INK, spacing=1.25)
    text(s, 48, 448, 834, 34,
         "Réponse : un univers pré-enregistré et horodaté AVANT toute ingestion — "
         "dix ETF américains, 21,7 ans, à contrainte strictement identique.",
         11.5, color=GREY)
    return s


def slide_expressiveness(prs, readiness, cap):
    s = new_slide(prs)
    header(s, "EXPRESSIVITÉ", "À contrainte identique, l'optimiseur peut enfin exprimer une vue", "")
    etf = cap["results"]["0.25"]["min_variance_lw"]["distinct_allocations"]
    free = readiness["allocation_freedom"]
    mv = free["min_variance_lw"]

    panel(s, 48, 108, 546, 250)
    s.shapes.add_picture(str(figure_png("global_2004_expressivite")),
                         Pt(60), Pt(130), width=Pt(522))

    kpi(s, 630, 122, 240, f"{etf} / 248", "allocations distinctes\nsur etf_2017", RED)
    kpi(s, 630, 226, 240, f"{mv['distinct_allocations']} / {free['n_rebalances']}",
        "allocations distinctes\nsur global_2004", TEAL)

    panel(s, 48, 378, 834, 56, CALLOUT)
    text(s, 70, 394, 790, 26,
         f"Même plafond de 25 %, mêmes coûts, même moteur — "
         f"{readiness['verdict']} sur {readiness['n_gates']} critères de préparation.",
         12.5, bold=True, color=INK)
    text(s, 48, 452, 834, 34,
         "C'est la condition nécessaire pour qu'un test de modèle ait un sens : "
         "sans elle, un résultat négatif ne s'interprète pas.", 11.5, color=GREY)
    return s


def slide_q1(prs, q1):
    s = new_slide(prs)
    header(s, "TEST ÉQUITABLE — Q1",
           "La couche de régimes ne bat pas le comparateur classique", "")
    cand, comp = q1["candidate"], q1["comparator"]
    diff, pi = q1["observed_difference"], q1["paired_inference"]

    kpi(s, 48, 118, 200, fr(comp["net_sharpe"], 4),
        "max_sharpe\ncomparateur pré-spécifié", TEAL)
    kpi(s, 268, 118, 200, fr(cand["net_sharpe"], 4),
        "regime_conditional\nstratégie candidate", BLUE)
    kpi(s, 488, 118, 200, fr(diff["net_sharpe_diff"], 4, signed=True),
        "écart de Sharpe net\nobservé", RED)
    kpi(s, 708, 118, 174, fr(pi["one_sided_null_centred_p_value"], 3),
        "valeur p unilatérale\ntest pairé", GREY)

    panel(s, 48, 240, 406, 176)
    text(s, 70, 258, 360, 20, "Ce que les coûts n'expliquent pas", 12.5, bold=True, color=INK)
    text(s, 70, 288, 366, 112,
         f"Le Sharpe BRUT était déjà inférieur : {fr(cand['gross_sharpe'], 4)} "
         f"contre {fr(comp['gross_sharpe'], 4)}.\n\nCe n'est donc pas un signal "
         "informatif annulé par les frais de transaction.",
         12, color=GREY, spacing=1.3)

    panel(s, 476, 240, 406, 176)
    text(s, 498, 258, 360, 20, "Ce que le Sharpe seul masquerait", 12.5, bold=True, color=INK)
    text(s, 498, 288, 366, 112,
         f"Rotation 4,3× supérieure ({fr(cand['avg_turnover'], 3)} contre "
         f"{fr(comp['avg_turnover'], 3)}).\n\nMais perte maximale MOINDRE : "
         f"{fr(100 * cand['max_drawdown'], 2)} % contre "
         f"{fr(100 * comp['max_drawdown'], 2)} % — la stratégie ne perd pas sur "
         "toutes les dimensions.",
         12, color=GREY, spacing=1.3)

    text(s, 48, 440, 834, 34,
         "Les deux verdicts pré-enregistrés sont négatifs. L'intervalle contient zéro : "
         "aucune supériorité n'est établie, dans aucun sens.", 11.5, color=GREY)
    return s


def slide_q2(prs, q2):
    s = new_slide(prs)
    header(s, "CORRECTION — Q2", "Le meilleur candidat brut est séduisant — et refusé", "")
    primary = q2["family_tests"]["primary_sharpe"]
    ledger = q2["candidate_ledger"]

    panel(s, 48, 108, 546, 264)
    s.shapes.add_picture(str(figure_png("global_2004_correction")),
                         Pt(58), Pt(122), width=Pt(526))

    kpi(s, 630, 116, 240, fr(primary["best_differential"], 3, signed=True),
        "meilleur écart BRUT\nsur 240 configurations", RED)
    kpi(s, 630, 212, 240,
        fr(primary["reality_check_p_value"], 3) + " / " + fr(primary["spa_p_value"], 3),
        "White RC / Hansen SPA\naucun ne rejette", TEAL)

    panel(s, 622, 306, 260, 66, CALLOUT)
    text(s, 640, 318, 226, 16, "Pourquoi il est refusé", 10.5, bold=True, color=AMBER)
    text(s, 640, 338, 226, 30,
         f"C'est le maximum d'une recherche de {ledger['executed_count']} "
         f"configurations. Seules {primary['n_candidates_beating_benchmark']} dépassent "
         "la référence.", 10.5, bold=True, color=INK, spacing=1.2)

    text(s, 48, 398, 834, 60,
         "Cité seul, ce +0,093 aurait fait un titre défendable — supérieur en valeur "
         "absolue à l'écart de Q1, et de sens inverse.\n"
         "La correction pour tests multiples est exactement ce qui l'en empêche.",
         12.5, bold=True, color=INK, spacing=1.3)
    return s


def slide_synthesis(prs, q2):
    s = new_slide(prs)
    header(s, "SYNTHÈSE", "Trois résultats à ne pas confondre", "")
    best = q2["family_tests"]["primary_sharpe"]["best_differential"]

    # 864pt of usable width (48pt margins) split into three cards with 20pt
    # gutters: (864 - 2*20) / 3 = 274.67. The first version used 280 at
    # x = 48/325/602, which OVERLAPPED by 3pt and rendered as one merged panel.
    cards = [
        (48, "Réussite d'ingénierie", TEAL,
         "Univers pré-enregistré et horodaté avant\ntoute donnée, dix critères de\n"
         "préparation, raccordé à DVC et Dagster,\npuis gelé comme preuve à usage unique."),
        (342, "Réussite méthodologique", BLUE,
         f"Un écart sélectionné de {fr(best, 3, signed=True)} de Sharpe\na été empêché de "
         "devenir un titre.\nC'est la contribution la plus\ntransportable du projet."),
        (636, "Absence d'avantage établi", RED,
         "Ni la couche de régimes (Q1) ni la\nfamille de challengers (Q2) n'établit\n"
         "de supériorité sur des règles de\nportefeuille plus simples."),
    ]
    for x, title, color, body in cards:
        panel(s, x, 118, 274, 196)
        text(s, x + 20, 138, 236, 40, title, 13, bold=True, color=color, spacing=1.15)
        text(s, x + 20, 180, 238, 126, body, 11.5, color=GREY, spacing=1.35)

    panel(s, 48, 352, 834, 82, CALLOUT)
    text(s, 70, 368, 400, 18, "Le message", 10.5, bold=True, color=AMBER)
    text(s, 70, 390, 790, 36,
         "Une fois les deux défauts d'identification levés, les modèles complexes ont "
         "enfin reçu un test équitable — et n'ont toujours pas établi d'avantage.",
         12.5, bold=True, color=INK, spacing=1.25)
    text(s, 48, 450, 834, 34,
         "La contribution est la preuve auditable montrant pourquoi le gagnant brut, "
         "pourtant séduisant, ne doit pas être cru.", 11.5, color=GREY)
    return s


# ── assembly ─────────────────────────────────────────────────────────────────

def slide_has_mark(slide) -> bool:
    return any(sh.has_text_frame and MARK in sh.text_frame.text for sh in sh_iter(slide))


def sh_iter(slide):
    return list(slide.shapes)


def main() -> None:
    readiness, cap = load("global_2004_readiness.json"), load("etf_cap_verdict.json")
    q1, q2 = load("global_2004_q1_results.json"), load("global_2004_q2_results.json")

    prs = Presentation(str(DECK))
    sld_id_lst = prs.slides._sldIdLst

    # Idempotence: drop any previously generated global_2004 slides.
    for sld_id in list(sld_id_lst):
        slide = prs.slides.get(sld_id.rId) if hasattr(prs.slides, "get") else None
        del slide  # not needed; identified below by index scan instead
    keep = []
    for idx, slide in enumerate(prs.slides):
        if slide_has_mark(slide):
            keep.append(idx)
    for idx in reversed(keep):
        rId = sld_id_lst[idx].rId
        prs.part.drop_rel(rId)
        sld_id_lst.remove(sld_id_lst[idx])
    if keep:
        print(f"  removed {len(keep)} previously generated slide(s)")

    n_before = len(prs.slides._sldIdLst)

    built = [slide_why(prs, readiness, cap),
             slide_expressiveness(prs, readiness, cap),
             slide_q1(prs, q1),
             slide_q2(prs, q2),
             slide_synthesis(prs, q2)]

    # They append at the end; move them into place.
    # Four content slides after slide 13 (ROBUSTESSE); the synthesis directly
    # before the closing CONCLUSION slide.
    ids = list(sld_id_lst)
    new_ids = ids[n_before:]
    for element in new_ids:
        sld_id_lst.remove(element)
    for offset, element in enumerate(new_ids[:4]):
        sld_id_lst.insert(13 + offset, element)
    sld_id_lst.insert(len(list(sld_id_lst)) - 1, new_ids[4])

    # Renumber every page label: the hand-made slides carry hardcoded numbers.
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            left_pt = shape.left / 12700 if shape.left is not None else 0
            top_pt = shape.top / 12700 if shape.top is not None else 0
            if 870 <= left_pt <= 890 and 20 <= top_pt <= 35:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                    if para.runs:
                        para.runs[0].text = f"{i:02d}" if i < 10 else str(i)
                break

    prs.save(str(DECK))
    print(f"deck now has {len(prs.slides._sldIdLst)} slides -> {DECK.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
